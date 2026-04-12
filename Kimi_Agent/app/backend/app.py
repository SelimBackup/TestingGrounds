"""
Backend API for Presentation to Excel Quotation Generator
Flask-based REST API for parsing PPTX/PDF files and generating Excel quotations
"""

import os
import re
import uuid
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pptx import Presentation
import pdfplumber
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Border, Side, Alignment
from openpyxl.utils import get_column_letter

# Determine if running in production (with built frontend) or development
FRONTEND_DIST = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'dist')

# Get allowed CORS origins from environment variable
CORS_ORIGINS = os.environ.get('CORS_ORIGINS', '*')
if CORS_ORIGINS != '*':
    CORS_ORIGINS = [origin.strip() for origin in CORS_ORIGINS.split(',')]

if os.path.exists(FRONTEND_DIST) and os.path.exists(os.path.join(FRONTEND_DIST, 'index.html')):
    # Production mode - serve static files
    app = Flask(__name__, static_folder=FRONTEND_DIST, static_url_path='')
    CORS(app, resources={r"/api/*": {"origins": CORS_ORIGINS}})
else:
    # Development mode - API only (standalone backend)
    app = Flask(__name__)
    CORS(app, resources={r"/api/*": {"origins": CORS_ORIGINS}})

# Configuration
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
OUTPUT_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'outputs')
ALLOWED_EXTENSIONS = {'pptx', 'pdf'}

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size


def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def clean_text(text):
    """Clean text by removing special characters and normalizing whitespace"""
    if not text:
        return ""
    text = ' '.join(text.split())
    return text.strip()


def extract_location(text, section=None):
    """Extract location from various formats"""
    # Pattern: -Location : Xxx or -Location: Xxx (stop at Type, Size, or newline)
    match = re.search(r'[-\s]*[Ll]ocation\s*:\s*([^\n]+?)(?=\s*Type:|\s*Size:|\n|$)', text)
    if match:
        return clean_text(match.group(1))
    
    # Pattern: Location: Xxx (without dash)
    match = re.search(r'[Ll]ocation\s*:\s*([^\n,]+)', text)
    if match:
        return clean_text(match.group(1))
    
    # Pattern: Word ending with period like "NewCairo."
    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        if re.match(r'^[A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+)?\.$', line):
            return line.rstrip('.')
    
    # Pattern: First line that looks like a location name (if no explicit Location field)
    for line in lines:
        line = line.strip()
        if any(line.lower().startswith(x) for x in ['size:', 'pixel:', 'description:', 'type:', '-']):
            continue
        if re.match(r'^[A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+)?$', line):
            return line
    
    return None


def extract_type(text):
    """Extract type/description from various formats"""
    # Pattern: Description :Xxx or Description: Xxx (stop at Size, Pixel, or newline)
    match = re.search(r'[Dd]escription\s*:\s*([^\n]+?)(?=\s*Size:|\s*Pixel:|\n|$)', text)
    if match:
        return clean_text(match.group(1))
    
    # Pattern: -Type: Xxx or Type: Xxx (stop at Size or newline)
    match = re.search(r'[-\s]*[Tt]ype\s*:\s*([^\n]+?)(?=\s*Size:|\n|$)', text)
    if match:
        return clean_text(match.group(1))
    
    # Pattern: standalone "Digital Screen" or similar
    match = re.search(r'^\s*([Dd]igital\s+[Ss]creen)\s*$', text, re.MULTILINE)
    if match:
        return match.group(1)
    
    return None


def extract_size_and_faces(text):
    """
    Extract all size entries with faces from text (handles multiple sizes per slide)
    Returns list of (width, length, faces) tuples
    
    Pattern examples:
    - Size: 8m x 17m –1 Face – 10 sec / 3 min
    - Size: 6m x 3m – 4 Faces- Pixel: 720 W X 360 H
    - Size: 8m x 16m –2 Face – 10 sec / 3 min
    - -Size : 10m x 19m – Digital Screen
    """
    entries = []
    
    # Pattern to match Size line with optional faces
    # Matches: [-]Size[: or :] Xm x Ym [separator] [N] Face(s)
    # Handles both "Size:" and "-Size :" formats with various dash types
    size_pattern = r'[-\s]*[Ss]ize\s*:\s*(\d+(?:\.\d+)?)\s*m?\s*[xX]\s*(\d+(?:\.\d+)?)\s*m?\s*(?:[–\-]\s*(\d+)?\s*[Ff]ace[s]?)?'
    
    matches = re.findall(size_pattern, text)
    
    for match in matches:
        try:
            width = float(match[0])
            length = float(match[1])
            # Extract faces, default to 1 if not specified
            faces = int(match[2]) if match[2] else 1
            entries.append((width, length, faces))
        except (ValueError, IndexError):
            continue
    
    return entries


def parse_slide_items(text, slide_num, section=None):
    """
    Parse slide text and extract all billboard items
    Returns list of data dictionaries (can be multiple per slide)
    """
    items = []
    
    # Get location
    location = extract_location(text, section)
    
    # Get type/description
    item_type = extract_type(text)
    
    # Get all size entries with faces
    size_entries = extract_size_and_faces(text)
    
    if size_entries:
        for width, length, faces in size_entries:
            items.append({
                'type': 'data_row',
                'slide_num': slide_num,
                'section': section,
                'location': location or section or "N/A",
                'item_type': item_type or "Digital Screen",
                'width': width,
                'length': length,
                'faces': faces,
                'meters': round(width * length, 2)
            })
    elif location or item_type:
        # Add item even without size (will have N/A for dimensions)
        items.append({
            'type': 'data_row',
            'slide_num': slide_num,
            'section': section,
            'location': location or section or "N/A",
            'item_type': item_type or "Digital Screen",
            'width': None,
            'length': None,
            'faces': 1,
            'meters': None
        })
    
    return items


def is_section_header_slide(text, slide_num):
    """
    Determine if a slide is a section header slide
    Section headers typically have:
    - Short text (just a location name)
    - No size information
    - No description
    """
    # Common non-section text patterns to exclude
    exclude_patterns = ['night vision', 'nightvision', 'digital screen']
    
    # If text is very short and has no size/pixel/description info, it's likely a header
    if len(text) < 50:
        lower_text = text.lower()
        if 'size:' not in lower_text and 'pixel:' not in lower_text and 'description:' not in lower_text:
            # Check if it looks like a location name
            lines = text.strip().split('\n')
            if len(lines) == 1:
                first_line = lines[0].strip()
                # Exclude common non-section patterns
                if any(pattern in first_line.lower() for pattern in exclude_patterns):
                    return False, None
                # Match patterns like "New Cairo", "Ring Road", "Nasr City"
                if re.match(r'^[A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+)?$', first_line):
                    return True, first_line
    return False, None


def has_meaningful_data(text):
    """Check if slide has meaningful data (size info) to extract"""
    lower_text = text.lower()
    # Must have size information to be considered a data slide
    # Match both "Size:" and "-Size :" formats
    if re.search(r'[-\s]*size\s*:', lower_text):
        return True
    return False


def extract_data_from_pptx(file_path):
    """
    Extract Location, Type, Width, Length, Faces, Meters from PPTX file
    Returns list of dicts with extracted data (no section headers)
    """
    prs = Presentation(file_path)
    extracted_data = []
    current_section = None
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Get all text from slide shapes
        slide_texts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_texts.append(text)
        
        full_text = '\n'.join(slide_texts)
        
        # Check if this is a section header slide
        is_header, header_name = is_section_header_slide(full_text, slide_num)
        
        if is_header:
            current_section = header_name
            # Don't add section headers to extracted data
            continue
        
        # Skip slides without meaningful data (like "Night Vision" slides)
        if not has_meaningful_data(full_text):
            continue
        
        # Extract data items from slide
        items = parse_slide_items(full_text, slide_num, current_section)
        extracted_data.extend(items)
    
    return extracted_data


def extract_data_from_pdf(file_path):
    """
    Extract Location, Type, Width, Length, Faces, Meters from PDF file
    Returns list of dicts with extracted data (no section headers)
    """
    extracted_data = []
    current_section = None
    
    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            
            # Check if this is a section header page
            is_header, header_name = is_section_header_slide(text, page_num)
            
            if is_header:
                current_section = header_name
                continue
            
            # Skip pages without meaningful data
            if not has_meaningful_data(text):
                continue
            
            # Extract data items
            items = parse_slide_items(text, page_num, current_section)
            extracted_data.extend(items)
    
    return extracted_data


def generate_excel_quotation(data_rows, unit_price, output_path):
    """
    Generate Excel quotation with formulas
    Columns: Location, Type, Width, Length, Faces, Meters, Unit Price, Print
    Print = Width × Length × Faces × Unit Price
    Unit Price is left empty for manual entry
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "Quotation"
    
    # Hide gridlines for professional look
    ws.sheet_view.showGridLines = False
    
    # Define styles
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=11)
    data_font = Font(size=10)
    currency_font = Font(size=10, color="0066CC")
    formula_font = Font(size=10, color="000000")  # Black for formula cells
    thin_border = Border(
        left=Side(style='thin', color='D0D0D0'),
        right=Side(style='thin', color='D0D0D0'),
        top=Side(style='thin', color='D0D0D0'),
        bottom=Side(style='thin', color='D0D0D0')
    )
    center_align = Alignment(horizontal='center', vertical='center')
    left_align = Alignment(horizontal='left', vertical='center')
    
    # Title
    ws.merge_cells('B2:I2')
    ws['B2'] = "QUOTATION"
    ws['B2'].font = Font(size=18, bold=True, color="1F4E79")
    ws['B2'].alignment = Alignment(horizontal='center', vertical='center')
    ws.row_dimensions[2].height = 30
    
    # Headers (row 4) - No more section headers in output
    headers = ['Location', 'Type', 'Width (m)', 'Length (m)', 'Faces', 'Meters', 'Unit Price', 'Print']
    header_row = 4
    
    for col_num, header in enumerate(headers, 2):  # Start from column B
        cell = ws.cell(row=header_row, column=col_num)
        cell.value = header
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = center_align
        cell.border = thin_border
    
    # Data rows
    current_row = header_row + 1
    data_start_row = current_row
    
    for item in data_rows:
        if item['type'] != 'data_row':
            continue  # Skip non-data rows
        
        # Location (Column B)
        cell = ws.cell(row=current_row, column=2)
        cell.value = item.get('location') or "N/A"
        cell.font = data_font
        cell.alignment = left_align
        cell.border = thin_border
        
        # Type (Column C)
        cell = ws.cell(row=current_row, column=3)
        cell.value = item.get('item_type') or "N/A"
        cell.font = data_font
        cell.alignment = left_align
        cell.border = thin_border
        
        # Width (Column D)
        cell = ws.cell(row=current_row, column=4)
        cell.value = item.get('width')
        if cell.value:
            cell.number_format = '0.00'
        cell.font = data_font
        cell.alignment = center_align
        cell.border = thin_border
        
        # Length (Column E)
        cell = ws.cell(row=current_row, column=5)
        cell.value = item.get('length')
        if cell.value:
            cell.number_format = '0.00'
        cell.font = data_font
        cell.alignment = center_align
        cell.border = thin_border
        
        # Faces (Column F)
        cell = ws.cell(row=current_row, column=6)
        cell.value = item.get('faces', 1)
        cell.font = data_font
        cell.alignment = center_align
        cell.border = thin_border
        
        # Meters (Column G) - Width × Length
        cell = ws.cell(row=current_row, column=7)
        width_col = get_column_letter(4)
        length_col = get_column_letter(5)
        cell.value = f"={width_col}{current_row}*{length_col}{current_row}"
        cell.number_format = '0.00'
        cell.font = formula_font
        cell.alignment = center_align
        cell.border = thin_border
        
        # Unit Price (Column H) - EMPTY for manual entry
        cell = ws.cell(row=current_row, column=8)
        cell.value = None  # Empty for manual entry
        cell.number_format = '$#,##0.00'
        cell.font = currency_font
        cell.alignment = center_align
        cell.border = thin_border
        
        # Print (Column I) = Width × Length × Faces × Unit Price
        cell = ws.cell(row=current_row, column=9)
        width_col = get_column_letter(4)
        length_col = get_column_letter(5)
        faces_col = get_column_letter(6)
        unit_price_col = get_column_letter(8)
        cell.value = f"={width_col}{current_row}*{length_col}{current_row}*{faces_col}{current_row}*{unit_price_col}{current_row}"
        cell.number_format = '$#,##0.00'
        cell.font = Font(size=10, bold=True)
        cell.alignment = center_align
        cell.border = thin_border
        
        current_row += 1
    
    data_end_row = current_row - 1
    
    # Grand Total for Print column
    if data_start_row <= data_end_row:
        total_row = current_row + 1
        ws.merge_cells(f'B{total_row}:H{total_row}')
        cell = ws.cell(row=total_row, column=2)
        cell.value = "GRAND TOTAL"
        cell.font = Font(bold=True, size=12, color="1F4E79")
        cell.alignment = Alignment(horizontal='right', vertical='center')
        cell.fill = PatternFill(start_color="E7E6E6", end_color="E7E6E6", fill_type="solid")
        
        cell = ws.cell(row=total_row, column=9)
        print_col = get_column_letter(9)
        total_formula = f"=SUM({print_col}{data_start_row}:{print_col}{data_end_row})"
        cell.value = total_formula
        cell.number_format = '$#,##0.00'
        cell.font = Font(bold=True, size=12, color="1F4E79")
        cell.alignment = center_align
        cell.fill = PatternFill(start_color="E7E6E6", end_color="E7E6E6", fill_type="solid")
    
    # Column widths
    column_widths = {
        'A': 2,   # Padding
        'B': 28,  # Location
        'C': 20,  # Type
        'D': 12,  # Width
        'E': 12,  # Length
        'F': 10,  # Faces
        'G': 12,  # Meters
        'H': 14,  # Unit Price
        'I': 16,  # Print
    }
    for col, width in column_widths.items():
        ws.column_dimensions[col].width = width
    
    # Instructions
    instruction_row = (total_row if data_start_row <= data_end_row else current_row) + 3
    ws.merge_cells(f'B{instruction_row}:I{instruction_row}')
    cell = ws.cell(row=instruction_row, column=2)
    cell.value = "Note: Enter Unit Price manually. Print column will auto-calculate (W × L × Faces × Unit Price)."
    cell.font = Font(size=9, italic=True, color="666666")
    
    wb.save(output_path)
    return output_path


@app.route('/api/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({'status': 'healthy', 'service': 'Presentation to Excel Quotation Generator'})


@app.route('/api/extract', methods=['POST'])
def extract_data():
    """Extract data from uploaded presentation file"""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type. Only PPTX and PDF files are allowed.'}), 400
    
    try:
        filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
        file.save(file_path)
        
        file_ext = filename.rsplit('.', 1)[1].lower()
        
        if file_ext == 'pptx':
            extracted_data = extract_data_from_pptx(file_path)
        elif file_ext == 'pdf':
            extracted_data = extract_data_from_pdf(file_path)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        os.remove(file_path)
        
        return jsonify({
            'success': True,
            'data': extracted_data,
            'count': len([d for d in extracted_data if d['type'] == 'data_row'])
        })
    
    except Exception as e:
        return jsonify({'error': f'Extraction failed: {str(e)}'}), 500


@app.route('/api/generate', methods=['POST'])
def generate_excel():
    """Generate Excel quotation from extracted data"""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    unit_price = request.form.get('unitPrice', '0')
    
    try:
        unit_price = float(unit_price)
    except ValueError:
        return jsonify({'error': 'Invalid unit price'}), 400
    
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type. Only PPTX and PDF files are allowed.'}), 400
    
    try:
        filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
        file.save(file_path)
        
        file_ext = filename.rsplit('.', 1)[1].lower()
        
        if file_ext == 'pptx':
            extracted_data = extract_data_from_pptx(file_path)
        elif file_ext == 'pdf':
            extracted_data = extract_data_from_pdf(file_path)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        os.remove(file_path)
        
        data_rows = [d for d in extracted_data if d['type'] == 'data_row']
        
        if not data_rows:
            return jsonify({'error': 'No data could be extracted from the file. Please check the file format.'}), 400
        
        output_filename = f"quotation_{unique_id}.xlsx"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        generate_excel_quotation(extracted_data, unit_price, output_path)
        
        return send_file(
            output_path,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=f"quotation_{filename.rsplit('.', 1)[0]}.xlsx"
        )
    
    except Exception as e:
        return jsonify({'error': f'Generation failed: {str(e)}'}), 500


@app.route('/', defaults={'path': ''})
@app.route('/<path:path>')
def serve_frontend(path):
    """Serve the React frontend in production mode"""
    if os.path.exists(FRONTEND_DIST) and os.path.exists(os.path.join(FRONTEND_DIST, 'index.html')):
        if path and os.path.exists(os.path.join(FRONTEND_DIST, path)):
            return send_file(os.path.join(FRONTEND_DIST, path))
        return send_file(os.path.join(FRONTEND_DIST, 'index.html'))
    return jsonify({'error': 'Frontend not built. Run npm run build first.'}), 404


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
